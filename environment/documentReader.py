#! python3
import PyPDF2
import docx
from pptx import Presentation


class pdfReader:

    def __init__(self, path:str):
        ''' Initiates the pdf reader object for the given file'''
        assert len(path) != 0
        self.pdfFileObj = open(path, 'rb')
        self.pdfReader = PyPDF2.PdfFileReader(self.pdfFileObj)
        self.pages = self.pdfReader.numPages

    def getText(self, start:int, end:int=None):
        ''' 
        Args: 
            start -> Starting page number (index starts from 1)
            end -> Endding page number
        Return:
            Dict with page number as key and the text content as value.
        '''
        pages = dict()
        for i in range(start-1, end):
            pageObj = self.pdfReader.getPage(i)
            pageContent = pageObj.extractText()
            pages[i] = pageContent
        
        return pages

class docxReader:

    def __init__(self, path:str):
        ''' Initiates the docx reader object for the given file'''
        assert len(path) != 0
        self.doc = docx.Document(path)

    def getText(self):
        ''' Returns the text content of the entire document '''    
        fullText = []
        for para in self.doc.paragraphs:
            fullText.append(para.text)
        return '\n'.join(fullText)

class pptReader:

    def __init__(self, path:str):
        ''' Initiates the ppt reader object for the given file'''
        assert len(path) != 0
        self.prs = Presentation(path)

    def getText(self):
        ''' Returns the text content of the entire ppt '''

        text_runs = []
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

        return text_runs

def main():
    lec1 = pdfReader('./Test Files/test.pdf')
    contents = lec1.getText(1,10)
    for k,v in contents.items():
        print(f'------------- Page {k} -------------')
        print(v)    
    print('\n======================================')

    lec2 = docxReader('./Test Files/test.docx')
    print('\n\nReading the word document...')
    print(lec2.getText())
    print('\n======================================')

    lec3 = pptReader('./Test Files/test.pptx')
    print(*lec3.getText())

if __name__ == "__main__":
    main()

    


